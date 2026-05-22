"""
Generate a professional, editable PPTX deck from AgentInvest markdown reports.
"""
import os
import re
import tempfile
import asyncio
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor
from typing import Any, Dict, List, Optional, Tuple


def _run_async_in_thread(coro):
    def _runner():
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        try:
            return loop.run_until_complete(coro)
        finally:
            loop.close()

    with ThreadPoolExecutor(max_workers=1) as executor:
        return executor.submit(_runner).result()


def _create_color_preserving_chart_html(chart_html: str, chartjs_src: str = None) -> str:
    return f"""<!DOCTYPE html>
<html style="background: white !important;">
<head>
    <meta charset="UTF-8">
    <meta name="color-scheme" content="light only">
    <script src="{chartjs_src or 'https://cdn.jsdelivr.net/npm/chart.js'}"></script>
    <style>
        * {{
            -webkit-print-color-adjust: exact !important;
            color-adjust: exact !important;
            print-color-adjust: exact !important;
            color-scheme: light !important;
            forced-color-adjust: none !important;
        }}
        body {{
            background: white !important;
            margin: 0 !important;
            padding: 20px !important;
        }}
        canvas {{
            background: transparent !important;
            max-width: 100% !important;
            height: auto !important;
        }}
    </style>
</head>
<body>
    {chart_html}
</body>
</html>"""


def _clean_markdown_text(text: str) -> str:
    """Convert markdown-ish content into cleaner plain text."""
    cleaned = re.sub(r"<[^>]+>", " ", text or "")
    cleaned = re.sub(r"```.*?```", " ", cleaned, flags=re.DOTALL)
    cleaned = re.sub(r"!\[[^\]]*\]\([^)]+\)", " ", cleaned)
    cleaned = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", cleaned)
    cleaned = re.sub(r"[*_`>#-]", " ", cleaned)
    cleaned = re.sub(r"\[(\d+)\]", "", cleaned)
    cleaned = re.sub(r"\s+", " ", cleaned)
    return cleaned.strip()


def _extract_executive_summary(report_markdown: str) -> str:
    pattern = r"##\s+Executive Summary\s*(.*?)(?=\n##\s+Table of Contents|\n##\s+[^\n]+|\Z)"
    match = re.search(pattern, report_markdown, flags=re.DOTALL)
    if not match:
        return ""
    return _clean_markdown_text(match.group(1))


def _extract_main_sections(report_markdown: str) -> List[Tuple[str, str]]:
    """Extract report sections from markdown body."""
    matches = list(re.finditer(r"^##\s+(.+)$", report_markdown, flags=re.MULTILINE))
    sections: List[Tuple[str, str]] = []
    ignored = {"Executive Summary", "Table of Contents", "References"}
    for idx, match in enumerate(matches):
        title = match.group(1).strip()
        if title in ignored:
            continue
        start = match.end()
        end = matches[idx + 1].start() if idx + 1 < len(matches) else len(report_markdown)
        body = _clean_markdown_text(report_markdown[start:end])
        if body:
            sections.append((title, body))
    return sections


def _extract_chart_blocks(report_markdown: str) -> List[str]:
    return re.findall(r"```html\s*\n(.*?)\n```", report_markdown or "", flags=re.DOTALL)



def _parse_chartjs_block(block: str, index: int) -> Optional[Dict[str, Any]]:
    """Parse a Chart.js config embedded in an AnalystIQ ```html block."""
    type_match = re.search(r"type:\s*['\"](\w+)['\"]", block)
    chart_type = (type_match.group(1) if type_match else "bar").lower()

    title = ""
    title_match = re.search(r"title:\s*\{[^}]*text:\s*['\"]([^'\"]+)['\"]", block, re.DOTALL)
    if title_match:
        title = title_match.group(1).strip()

    labels_match = re.search(r"labels:\s*\[(.*?)\]", block, re.DOTALL)
    categories: List[str] = []
    if labels_match:
        categories = re.findall(r"['\"]([^'\"]+)['\"]", labels_match.group(1))

    series_label = "Series"
    values: List[float] = []
    dataset_match = re.search(r"datasets:\s*\[\s*\{([\s\S]*?)\}\s*\]", block)
    if dataset_match:
        dataset_body = dataset_match.group(1)
        label_match = re.search(r"label:\s*['\"]([^'\"]+)['\"]", dataset_body)
        if label_match:
            series_label = label_match.group(1).strip()
        data_match = re.search(r"data:\s*\[([\s\S]*?)\]", dataset_body)
        if data_match:
            for token in re.findall(r"-?\d[\d,]*(?:\.\d+)?", data_match.group(1)):
                try:
                    values.append(float(token.replace(",", "")))
                except ValueError:
                    continue

    if not categories or not values:
        return None

    ppt_type = "BAR"
    if chart_type in {"line", "area"}:
        ppt_type = "LINE"
    elif chart_type in {"pie", "doughnut", "donut"}:
        ppt_type = "PIE"

    pair_count = min(len(categories), len(values))
    return {
        "index": index,
        "source_type": chart_type,
        "ppt_type": ppt_type,
        "title": title or f"Chart {index + 1}",
        "categories": categories[:pair_count],
        "values": values[:pair_count],
        "series_label": series_label,
    }


def extract_chart_specs_from_markdown(markdown_content: str) -> List[Dict[str, Any]]:
    """Extract structured chart metadata from Chart.js HTML blocks in report markdown."""
    charts: List[Dict[str, Any]] = []
    for idx, block in enumerate(_extract_chart_blocks(markdown_content)):
        spec = _parse_chartjs_block(block, idx)
        if spec:
            charts.append(spec)
    return charts


async def _render_chart_images(chart_blocks: List[str], output_dir: str, chartjs_src: str = None) -> List[str]:
    if not chart_blocks:
        return []
    try:
        from playwright.async_api import async_playwright, TimeoutError as PlaywrightTimeoutError
    except ImportError:
        return []

    image_paths: List[str] = []
    async with async_playwright() as p:
        browser = await p.chromium.launch(
            headless=True,
            args=["--no-sandbox", "--disable-dev-shm-usage", "--force-color-profile=srgb"],
        )
        try:
            for idx, chart_html in enumerate(chart_blocks):
                page = await browser.new_page(viewport={"width": 1280, "height": 720}, color_scheme="light")
                chart_doc = _create_color_preserving_chart_html(chart_html, chartjs_src)
                await page.set_content(chart_doc)
                await page.wait_for_function("typeof Chart !== 'undefined'", timeout=15000)
                try:
                    await page.wait_for_selector("canvas", timeout=15000)
                except PlaywrightTimeoutError:
                    # Some generated blocks may not produce a canvas; continue gracefully.
                    pass
                await page.wait_for_timeout(2500)
                path = os.path.join(output_dir, f"chart_{idx}.png")
                await page.screenshot(path=path, type="png", full_page=False)
                await page.close()
                if os.path.exists(path) and os.path.getsize(path) > 0:
                    image_paths.append(path)
        finally:
            await browser.close()
    return image_paths


def _truncate(text: str, max_chars: int = 650) -> str:
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 3].rstrip() + "..."


def _compress_bullet(text: str, max_words: int = 14, max_chars: int = 110) -> str:
    clean = _clean_markdown_text(text)
    words = clean.split()
    short = " ".join(words[:max_words])
    if len(short) > max_chars:
        short = short[: max_chars - 3].rstrip() + "..."
    return short


def _fallback_slide_plan(
    report_markdown: str, key_points: List[str], executive_summary: str
) -> List[Dict[str, Any]]:
    sections = _extract_main_sections(report_markdown)
    plan: List[Dict[str, Any]] = []

    if executive_summary.strip():
        summary_chunks = [s.strip() for s in re.split(r"(?<=[.!?])\s+", executive_summary) if s.strip()]
        summary_bullets = [_compress_bullet(chunk) for chunk in summary_chunks[:4] if chunk]
        if summary_bullets:
            plan.append(
                {
                    "title": "Investment Thesis",
                    "subtitle": "High-level view",
                    "bullets": summary_bullets[:4],
                    "use_chart": True,
                }
            )

    if key_points:
        plan.append(
            {
                "title": "Key Decision Highlights",
                "subtitle": "Most material findings",
                "bullets": [_compress_bullet(point) for point in key_points[:5]],
                "use_chart": False,
            }
        )

    for idx, (title, content) in enumerate(sections[:4]):
        sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", content) if s.strip()]
        bullets = [_compress_bullet(s) for s in sentences[:4] if s]
        if bullets:
            plan.append(
                {
                    "title": title,
                    "subtitle": "",
                    "bullets": bullets,
                    "use_chart": idx % 2 == 0,
                }
            )
    return plan[:8]


def _normalize_slide_plan(
    slide_plan: Optional[List[Dict[str, Any]]],
    report_markdown: str,
    key_points: List[str],
    executive_summary: str,
) -> List[Dict[str, Any]]:
    source = slide_plan or _fallback_slide_plan(report_markdown, key_points, executive_summary)
    normalized: List[Dict[str, Any]] = []
    for item in source:
        if not isinstance(item, dict):
            continue
        title = str(item.get("title", "")).strip()[:80]
        if not title:
            continue
        subtitle = str(item.get("subtitle", "")).strip()[:120]
        bullets_value = item.get("bullets", [])
        bullets: List[str] = []
        if isinstance(bullets_value, list):
            for bullet in bullets_value[:5]:
                compressed = _compress_bullet(str(bullet))
                if compressed:
                    bullets.append(compressed)
        if len(bullets) < 3:
            continue
        normalized.append(
            {
                "title": title,
                "subtitle": subtitle,
                "bullets": bullets[:5],
                "use_chart": bool(item.get("use_chart", False)),
            }
        )
    return normalized[:8]


def build_professional_pptx(
    *,
    report_markdown: str,
    output_path: str,
    company_name: str,
    ticker: str,
    key_points: List[str],
    executive_summary: str = "",
    chartjs_src: str = None,
    slide_plan: Optional[List[Dict[str, Any]]] = None,
    visual_deck_spec: Optional[Dict[str, Any]] = None,
    style_profile: str = "Institutional Light",
) -> str:
    """Build an editable, visual investment presentation with at most 10 slides."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    style_presets = {
        "Institutional Light": {
            "brand_navy": RGBColor(15, 23, 42),
            "accent_blue": RGBColor(37, 99, 235),
            "text_slate": RGBColor(51, 65, 85),
            "muted_gray": RGBColor(100, 116, 139),
            "bg_color": RGBColor(247, 250, 255),
            "card_color": RGBColor(255, 255, 255),
            "line_color": RGBColor(219, 234, 254),
        },
        "Executive Dark": {
            "brand_navy": RGBColor(240, 246, 255),
            "accent_blue": RGBColor(96, 165, 250),
            "text_slate": RGBColor(226, 232, 240),
            "muted_gray": RGBColor(148, 163, 184),
            "bg_color": RGBColor(15, 23, 42),
            "card_color": RGBColor(30, 41, 59),
            "line_color": RGBColor(51, 65, 85),
        },
        "Minimal Clean": {
            "brand_navy": RGBColor(17, 24, 39),
            "accent_blue": RGBColor(14, 116, 144),
            "text_slate": RGBColor(55, 65, 81),
            "muted_gray": RGBColor(107, 114, 128),
            "bg_color": RGBColor(255, 255, 255),
            "card_color": RGBColor(255, 255, 255),
            "line_color": RGBColor(229, 231, 235),
        },
    }
    style = style_presets.get(style_profile, style_presets["Institutional Light"])

    brand_navy = style["brand_navy"]
    accent_blue = style["accent_blue"]
    text_slate = style["text_slate"]
    muted_gray = style["muted_gray"]
    bg_color = style["bg_color"]
    card_color = style["card_color"]
    line_color = style["line_color"]
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    margin = Inches(0.55)
    footer_y = Inches(7.02)
    content_bottom = Inches(6.78)
    summary_text = executive_summary.strip() or _extract_executive_summary(report_markdown)

    def plain(value: Any, max_chars: int = 140) -> str:
        raw = str(value or "")
        replacements = {
            "\u2018": "'",
            "\u2019": "'",
            "\u201c": '"',
            "\u201d": '"',
            "\u2013": "-",
            "\u2014": "-",
            "\ufffd": "'",
        }
        for old, new in replacements.items():
            raw = raw.replace(old, new)
        cleaned = _clean_markdown_text(raw)
        cleaned = re.sub(r"^(Executive Summary|Investment Recommendation|Confidence Level):\s*", "", cleaned, flags=re.I)
        cleaned = re.sub(r"\b(Investment Recommendation|Confidence Level):\s*[^.]+\.?\s*", "", cleaned, flags=re.I)
        if max_chars and len(cleaned) > max_chars:
            cleaned = cleaned[: max_chars - 3].rstrip() + "..."
        return cleaned

    def clean_heading(value: Any, max_chars: int = 78) -> str:
        heading = plain(value, max_chars=160)
        heading = re.sub(r"^\d+[\.\)]\s*", "", heading).strip()
        heading = re.sub(r"\s+", " ", heading)
        if len(heading) > max_chars:
            heading = heading[: max_chars - 3].rstrip() + "..."
        return heading

    def make_legacy_visual_spec() -> Dict[str, Any]:
        normalized_plan = _normalize_slide_plan(slide_plan, report_markdown, key_points, summary_text)
        slides: List[Dict[str, Any]] = []
        for idx, spec in enumerate(normalized_plan):
            slides.append(
                {
                    "layout_type": "chart_focus" if spec.get("use_chart") else "two_column",
                    "section_label": clean_heading(spec.get("subtitle") or f"Section {idx + 1}", max_chars=32),
                    "headline": clean_heading(spec.get("title", "Investment Update")),
                    "takeaway": plain((spec.get("bullets") or [""])[0], 115),
                    "bullets": (spec.get("bullets") or [])[1:5],
                    "metrics": [],
                    "chart_ref": idx if spec.get("use_chart") else None,
                    "visual_emphasis": "",
                    "speaker_notes": "",
                }
            )
        if not slides:
            slides = _fallback_visual_spec()["slides"]
        return {
            "deck_title": f"{company_name} Investment Committee Deck",
            "subtitle": f"{ticker} | Generated by AgentInvest",
            "investment_thesis": plain(summary_text, 155),
            "recommendation": "Validate thesis, risks, and sizing before committee action.",
            "slides": slides[:9],
        }

    def _fallback_visual_spec() -> Dict[str, Any]:
        sections = _extract_main_sections(report_markdown)
        thesis_bullets = [plain(point, max_chars=95) for point in key_points[:3]]
        slides = [
            {
                "layout_type": "thesis",
                "section_label": "Thesis",
                "headline": "Investment thesis and key decision points",
                "takeaway": plain(summary_text, 145),
                "bullets": thesis_bullets,
                "metrics": [],
                "chart_ref": 0,
            }
        ]
        for idx, (title, content) in enumerate(sections[:5]):
            sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", content) if s.strip()]
            bullets = [plain(sentence, max_chars=92) for sentence in sentences[:4]]
            heading = clean_heading(title)
            slides.append(
                {
                    "layout_type": "chart_focus" if idx % 2 == 0 else "two_column",
                    "section_label": heading[:32],
                    "headline": heading,
                    "takeaway": bullets[0] if bullets else "",
                    "bullets": bullets[1:] if len(bullets) > 1 else bullets,
                    "metrics": [],
                    "chart_ref": idx if idx % 2 == 0 else None,
                }
            )
        slides.append(
            {
                "layout_type": "closing_recommendation",
                "section_label": "Recommendation",
                "headline": "Committee actions and monitoring plan",
                "takeaway": "Convert the report into sizing, risk, and catalyst decisions.",
                "bullets": [
                    "Validate assumptions against internal model and consensus.",
                    "Pressure-test downside scenarios before position sizing.",
                    "Define buy, hold, trim, or watchlist decision criteria.",
                ],
                "metrics": [],
                "chart_ref": None,
            }
        )
        return {
            "deck_title": f"{company_name} Investment Committee Deck",
            "subtitle": f"{ticker} | Generated by AgentInvest",
            "investment_thesis": plain(summary_text, 155),
            "recommendation": "Validate thesis, risks, and sizing before committee action.",
            "slides": slides[:9],
        }

    deck_spec = visual_deck_spec or make_legacy_visual_spec()
    if not isinstance(deck_spec, dict) or not deck_spec.get("slides"):
        deck_spec = _fallback_visual_spec()

    def blank_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color
        return slide

    def add_text(
        slide,
        text: str,
        left,
        top,
        width,
        height,
        *,
        size: int = 18,
        color=None,
        bold: bool = False,
        align=None,
        all_caps: bool = False,
    ):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.03)
        tf.margin_right = Inches(0.03)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        paragraph = tf.paragraphs[0]
        paragraph.text = plain(text.upper() if all_caps else text, max_chars=420)
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color or text_slate
        if align is not None:
            paragraph.alignment = align
        return box

    def add_card(slide, left, top, width, height, fill=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill or card_color
        shape.line.color.rgb = line_color
        return shape

    def add_footer(slide, slide_no: int):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, margin, Inches(6.93), slide_w - Inches(1.1), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = line_color
        line.line.fill.background()
        add_text(slide, "AnalystIQ", margin, footer_y, Inches(2.2), Inches(0.25), size=8, color=muted_gray)
        add_text(
            slide,
            f"{company_name} ({ticker})",
            Inches(5.2),
            footer_y,
            Inches(2.9),
            Inches(0.25),
            size=8,
            color=muted_gray,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            str(slide_no),
            Inches(12.2),
            footer_y,
            Inches(0.55),
            Inches(0.25),
            size=8,
            color=muted_gray,
            align=PP_ALIGN.RIGHT,
        )

    def add_header(slide, spec: Dict[str, Any], slide_no: int):
        add_text(
            slide,
            clean_heading(spec.get("section_label") or "Investment View", max_chars=32),
            margin,
            Inches(0.22),
            Inches(3.0),
            Inches(0.28),
            size=8,
            color=accent_blue,
            bold=True,
            all_caps=True,
        )
        add_text(
            slide,
            clean_heading(spec.get("headline") or "Investment Update"),
            margin,
            Inches(0.58),
            Inches(11.8),
            Inches(0.72),
            size=27,
            color=brand_navy,
            bold=True,
        )
        add_footer(slide, slide_no)

    def add_bullets(slide, bullets: List[str], left, top, width, height, *, size: int = 16, max_items: int = 4):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.06)
        clean_bullets = [plain(bullet, 105) for bullet in bullets if plain(bullet, 105)][:max_items]
        if not clean_bullets:
            return box
        for idx, bullet in enumerate(clean_bullets):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(size)
            p.font.color.rgb = text_slate
            p.space_after = Pt(8)
        return box

    def add_metrics(slide, metrics: List[Dict[str, Any]], left, top, width, height):
        if not metrics:
            return
        count = min(len(metrics), 4)
        gap = Inches(0.12)
        card_w = int((width - gap * (count - 1)) / count)
        for idx, metric in enumerate(metrics[:count]):
            x = left + idx * (card_w + gap)
            add_card(slide, x, top, card_w, height)
            add_text(slide, metric.get("label", "Metric"), x + Inches(0.15), top + Inches(0.14), card_w - Inches(0.3), Inches(0.28), size=8, color=muted_gray, bold=True, all_caps=True)
            add_text(slide, metric.get("value", ""), x + Inches(0.15), top + Inches(0.42), card_w - Inches(0.3), Inches(0.42), size=18, color=brand_navy, bold=True)
            if metric.get("delta"):
                add_text(slide, metric.get("delta", ""), x + Inches(0.15), top + Inches(0.88), card_w - Inches(0.3), Inches(0.32), size=9, color=accent_blue)

    def png_size(path: str) -> Optional[Tuple[int, int]]:
        try:
            with open(path, "rb") as image_file:
                header = image_file.read(24)
            if header.startswith(b"\x89PNG\r\n\x1a\n") and len(header) >= 24:
                width = int.from_bytes(header[16:20], "big")
                height = int.from_bytes(header[20:24], "big")
                if width > 0 and height > 0:
                    return width, height
        except OSError:
            return None
        return None

    def add_fit_picture(slide, image_path: str, left, top, width, height):
        size = png_size(image_path)
        aspect = (size[0] / size[1]) if size else (16 / 9)
        box_ratio = width / height
        if aspect >= box_ratio:
            pic_w = width
            pic_h = int(width / aspect)
        else:
            pic_h = height
            pic_w = int(height * aspect)
        pic_left = left + int((width - pic_w) / 2)
        pic_top = top + int((height - pic_h) / 2)
        return slide.shapes.add_picture(image_path, pic_left, pic_top, width=pic_w, height=pic_h)

    def add_native_chart(slide, chart_spec: Dict[str, Any], left, top, width, height):
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

        if not chart_spec or not chart_spec.get("categories") or not chart_spec.get("values"):
            return None
        chart_data = CategoryChartData()
        chart_data.categories = chart_spec["categories"]
        chart_data.add_series(chart_spec.get("series_label") or "Series", tuple(chart_spec["values"]))
        ppt_type = str(chart_spec.get("ppt_type") or "BAR").upper()
        chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
        if ppt_type == "LINE":
            chart_type = XL_CHART_TYPE.LINE
        elif ppt_type == "PIE":
            chart_type = XL_CHART_TYPE.PIE
        chart_shape = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
        chart = chart_shape.chart
        chart.has_title = True
        chart.chart_title.text_frame.text = plain(chart_spec.get("title") or "Chart", max_chars=80)
        chart.has_legend = ppt_type != "PIE"
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        if chart.plots and ppt_type in {"BAR", "LINE"}:
            chart.plots[0].has_data_labels = True
        return chart_shape

    def resolve_chart_spec(spec: Dict[str, Any], chart_specs: List[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
        chart_ref = spec.get("chart_ref")
        if isinstance(chart_ref, int) and 0 <= chart_ref < len(chart_specs):
            return chart_specs[chart_ref]
        return None

    def resolve_chart_path(spec: Dict[str, Any], chart_paths: List[str], used: set) -> Optional[str]:
        if not chart_paths:
            return None
        chart_ref = spec.get("chart_ref")
        if isinstance(chart_ref, int) and 0 <= chart_ref < len(chart_paths):
            used.add(chart_ref)
            return chart_paths[chart_ref]
        for idx, path in enumerate(chart_paths):
            if idx not in used:
                used.add(idx)
                return path
        return None

    def render_title_slide():
        slide = blank_slide()
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), slide_h)
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_blue
        accent_bar.line.fill.background()
        add_text(slide, "ANALYSTIQ", margin, Inches(0.35), Inches(2.0), Inches(0.25), size=9, color=accent_blue, bold=True)
        add_text(
            slide,
            deck_spec.get("deck_title") or f"{company_name} Investment Committee Deck",
            Inches(1.5),
            Inches(2.2),
            Inches(10.333),
            Inches(1.6),
            size=46,
            color=brand_navy,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        subtitle = deck_spec.get("subtitle") or f"{ticker} | Generated by AnalystIQ"
        add_text(
            slide,
            subtitle,
            Inches(1.5),
            Inches(3.95),
            Inches(10.333),
            Inches(0.55),
            size=20,
            color=accent_blue,
            align=PP_ALIGN.CENTER,
        )
        thesis = plain(deck_spec.get("investment_thesis") or summary_text, 155)
        if thesis:
            add_card(slide, Inches(2.4), Inches(4.85), Inches(8.5), Inches(1.35))
            add_text(slide, thesis, Inches(2.75), Inches(5.15), Inches(7.8), Inches(0.85), size=14, color=text_slate, align=PP_ALIGN.CENTER)
        add_text(slide, datetime.now().strftime("%Y-%m-%d"), margin, Inches(6.6), Inches(2.0), Inches(0.3), size=10, color=muted_gray)
        add_footer(slide, 1)

    def render_big_stat_slide(spec: Dict[str, Any], slide_no: int):
        slide = blank_slide()
        add_header(slide, spec, slide_no)
        stat_value = plain(spec.get("stat_number") or spec.get("takeaway") or "—", max_chars=18)
        stat_label = plain(spec.get("stat_label") or spec.get("subheading") or spec.get("headline") or "", max_chars=120)
        add_text(slide, stat_value, Inches(1.0), Inches(2.4), Inches(5.0), Inches(2.5), size=88, color=accent_blue, bold=True)
        add_text(slide, stat_label, Inches(6.2), Inches(3.1), Inches(6.0), Inches(2.0), size=24, color=brand_navy)
        bullets = spec.get("bullets") or []
        if bullets:
            add_bullets(slide, bullets, Inches(6.2), Inches(4.35), Inches(6.0), Inches(2.0), size=14, max_items=2)

    def render_three_column_slide(spec: Dict[str, Any], slide_no: int):
        slide = blank_slide()
        add_header(slide, spec, slide_no)
        bullets = (spec.get("bullets") or [])[:3]
        for idx, point in enumerate(bullets):
            left = Inches(1.0) + (idx * Inches(3.9))
            add_card(slide, left, Inches(2.6), Inches(3.5), Inches(4.0))
            add_text(slide, f"// 0{idx + 1}", left + Inches(0.28), Inches(2.95), Inches(2.9), Inches(0.3), size=14, color=accent_blue, bold=True)
            add_text(slide, plain(point, 120), left + Inches(0.28), Inches(3.45), Inches(2.9), Inches(2.8), size=16, color=text_slate)

    def render_text_and_image_slide(
        spec: Dict[str, Any],
        slide_no: int,
        chart_path: Optional[str],
        chart_spec: Optional[Dict[str, Any]],
    ):
        slide = blank_slide()
        add_header(slide, spec, slide_no)
        add_bullets(slide, spec.get("bullets", []), Inches(1.0), Inches(2.2), Inches(5.5), Inches(4.5), size=16, max_items=3)
        panel_left = Inches(7.0)
        panel_top = Inches(2.2)
        panel_w = Inches(5.5)
        panel_h = Inches(4.5)
        add_card(slide, panel_left, panel_top, panel_w, panel_h)
        if chart_path:
            add_fit_picture(slide, chart_path, panel_left + Inches(0.2), panel_top + Inches(0.2), panel_w - Inches(0.4), panel_h - Inches(0.4))
        elif chart_spec:
            add_native_chart(slide, chart_spec, panel_left + Inches(0.25), panel_top + Inches(0.35), panel_w - Inches(0.5), panel_h - Inches(0.55))
        else:
            emphasis = plain(spec.get("visual_emphasis") or spec.get("subheading") or "Key investment implication", max_chars=90)
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, panel_left + Inches(0.35), panel_top + Inches(0.35), Inches(0.12), panel_h - Inches(0.7))
            accent.fill.solid()
            accent.fill.fore_color.rgb = accent_blue
            accent.line.fill.background()
            add_text(slide, emphasis, panel_left + Inches(0.65), panel_top + Inches(1.5), panel_w - Inches(1.0), Inches(2.2), size=22, color=brand_navy, bold=True)

    def render_hero_slide(spec: Dict[str, Any], slide_no: int):
        slide = blank_slide()
        add_header(slide, spec, slide_no)
        add_card(slide, margin, Inches(1.55), Inches(12.1), Inches(4.85))
        add_text(slide, plain(spec.get("takeaway") or "", 145), Inches(0.95), Inches(2.05), Inches(11.2), Inches(1.0), size=24, color=brand_navy, bold=True, align=PP_ALIGN.CENTER)
        add_bullets(slide, spec.get("bullets", []), Inches(2.2), Inches(3.35), Inches(8.9), Inches(2.2), size=17, max_items=3)

    def render_risk_matrix_slide(spec: Dict[str, Any], slide_no: int, chart_path: Optional[str]):
        slide = blank_slide()
        add_header(slide, spec, slide_no)
        add_card(slide, Inches(0.72), Inches(1.48), Inches(5.65), Inches(5.05))
        add_text(slide, "Risk Lens", Inches(1.0), Inches(1.75), Inches(2.5), Inches(0.25), size=9, color=accent_blue, bold=True, all_caps=True)
        add_bullets(slide, spec.get("bullets", []), Inches(1.0), Inches(2.2), Inches(5.05), Inches(3.8), size=15, max_items=3)
        right = Inches(6.72)
        top = Inches(1.48)
        add_card(slide, right, top, Inches(5.9), Inches(5.05))
        add_text(slide, plain(spec.get("takeaway") or "", 125), right + Inches(0.35), top + Inches(0.45), Inches(5.1), Inches(1.2), size=20, color=brand_navy, bold=True)
        if chart_path:
            add_fit_picture(slide, chart_path, right + Inches(0.22), top + Inches(1.85), Inches(5.46), Inches(2.85))
        else:
            add_text(slide, "Mitigation focus", right + Inches(0.35), top + Inches(2.2), Inches(5.1), Inches(0.25), size=9, color=accent_blue, bold=True, all_caps=True)
            add_text(slide, plain(spec.get("visual_emphasis") or "Prioritize high-impact risks with explicit monitoring triggers.", 120), right + Inches(0.35), top + Inches(2.65), Inches(5.1), Inches(2.2), size=15, color=text_slate)


    def render_thesis_slide(spec: Dict[str, Any], slide_no: int, chart_path: Optional[str]):
        slide = blank_slide()
        add_header(slide, spec, slide_no)
        add_card(slide, margin, Inches(1.48), Inches(5.3), Inches(4.95))
        add_text(slide, "Core Takeaway", Inches(0.85), Inches(1.78), Inches(4.7), Inches(0.3), size=10, color=accent_blue, bold=True, all_caps=True)
        add_text(slide, plain(spec.get("takeaway") or deck_spec.get("investment_thesis"), 145), Inches(0.85), Inches(2.15), Inches(4.65), Inches(1.1), size=18, color=brand_navy, bold=True)
        add_bullets(slide, spec.get("bullets", []), Inches(0.85), Inches(3.55), Inches(4.65), Inches(2.35), size=13, max_items=3)
        if chart_path:
            add_card(slide, Inches(6.15), Inches(1.48), Inches(6.55), Inches(4.95))
            add_fit_picture(slide, chart_path, Inches(6.35), Inches(1.75), Inches(6.15), Inches(4.35))
        else:
            add_metrics(slide, spec.get("metrics", []), Inches(6.15), Inches(1.65), Inches(6.1), Inches(1.35))

    def render_chart_focus_slide(
        spec: Dict[str, Any],
        slide_no: int,
        chart_path: Optional[str],
        chart_spec: Optional[Dict[str, Any]] = None,
    ):
        slide = blank_slide()
        add_header(slide, spec, slide_no)
        chart_left = Inches(0.72)
        chart_top = Inches(1.42)
        chart_w = Inches(8.25)
        chart_h = Inches(5.15)
        add_card(slide, chart_left, chart_top, chart_w, chart_h)
        if chart_path:
            add_fit_picture(slide, chart_path, chart_left + Inches(0.2), chart_top + Inches(0.2), chart_w - Inches(0.4), chart_h - Inches(0.4))
        elif chart_spec:
            add_native_chart(slide, chart_spec, chart_left + Inches(0.25), chart_top + Inches(0.25), chart_w - Inches(0.5), chart_h - Inches(0.5))
        else:
            add_text(slide, "No chart available for this slide", chart_left + Inches(0.4), chart_top + Inches(2.2), chart_w - Inches(0.8), Inches(0.4), size=16, color=muted_gray, align=PP_ALIGN.CENTER)
        panel_left = Inches(9.25)
        add_card(slide, panel_left, chart_top, Inches(3.35), chart_h)
        context = plain(spec.get("subheading") or spec.get("takeaway") or "", 115)
        add_text(slide, "Insight", panel_left + Inches(0.28), chart_top + Inches(0.28), Inches(2.75), Inches(0.28), size=9, color=accent_blue, bold=True, all_caps=True)
        add_text(slide, context, panel_left + Inches(0.28), chart_top + Inches(0.68), Inches(2.75), Inches(1.25), size=15, color=brand_navy, bold=True)
        add_bullets(slide, spec.get("bullets", []), panel_left + Inches(0.28), chart_top + Inches(2.1), Inches(2.75), Inches(2.6), size=12, max_items=3)

    def render_two_column_slide(spec: Dict[str, Any], slide_no: int, chart_path: Optional[str]):
        slide = blank_slide()
        add_header(slide, spec, slide_no)
        left = Inches(0.72)
        right = Inches(6.72)
        top = Inches(1.48)
        add_card(slide, left, top, Inches(5.65), Inches(5.05))
        add_text(slide, plain(spec.get("takeaway", ""), 125), left + Inches(0.28), top + Inches(0.35), Inches(5.05), Inches(1.0), size=18, color=brand_navy, bold=True)
        add_bullets(slide, spec.get("bullets", []), left + Inches(0.28), top + Inches(1.6), Inches(5.05), Inches(3.0), size=15, max_items=4)
        add_card(slide, right, top, Inches(5.9), Inches(5.05))
        if chart_path:
            add_fit_picture(slide, chart_path, right + Inches(0.22), top + Inches(0.32), Inches(5.46), Inches(4.35))
        elif spec.get("metrics"):
            add_metrics(slide, spec.get("metrics", []), right + Inches(0.3), top + Inches(0.55), Inches(5.25), Inches(1.35))
        else:
            add_text(slide, spec.get("visual_emphasis") or "Implication for the investment case", right + Inches(0.35), top + Inches(1.8), Inches(5.1), Inches(0.8), size=22, color=accent_blue, bold=True, align=PP_ALIGN.CENTER)

    def render_metrics_slide(spec: Dict[str, Any], slide_no: int, chart_path: Optional[str]):
        slide = blank_slide()
        add_header(slide, spec, slide_no)
        add_metrics(slide, spec.get("metrics", []), margin, Inches(1.48), Inches(12.1), Inches(1.35))
        add_card(slide, margin, Inches(3.1), Inches(5.1), Inches(3.3))
        add_text(slide, plain(spec.get("takeaway", ""), 125), Inches(0.85), Inches(3.45), Inches(4.45), Inches(0.9), size=17, color=brand_navy, bold=True)
        add_bullets(slide, spec.get("bullets", []), Inches(0.85), Inches(4.55), Inches(4.45), Inches(1.5), size=13, max_items=3)
        if chart_path:
            add_card(slide, Inches(6.0), Inches(3.1), Inches(6.35), Inches(3.3))
            add_fit_picture(slide, chart_path, Inches(6.25), Inches(3.3), Inches(5.85), Inches(2.85))

    def render_closing_slide(spec: Dict[str, Any], slide_no: int):
        slide = blank_slide()
        add_header(slide, spec, slide_no)
        add_card(slide, margin, Inches(1.55), Inches(12.1), Inches(4.85))
        add_text(slide, plain(spec.get("takeaway") or deck_spec.get("recommendation"), 145), Inches(0.95), Inches(1.95), Inches(11.2), Inches(0.8), size=23, color=brand_navy, bold=True, align=PP_ALIGN.CENTER)
        add_bullets(slide, spec.get("bullets", []), Inches(2.2), Inches(3.15), Inches(8.9), Inches(2.2), size=17, max_items=4)

    render_title_slide()
    chart_blocks = _extract_chart_blocks(report_markdown)
    chart_specs = extract_chart_specs_from_markdown(report_markdown)
    with tempfile.TemporaryDirectory() as temp_dir:
        chart_paths = _run_async_in_thread(_render_chart_images(chart_blocks, temp_dir, chartjs_src))
        used_chart_refs: set = set()
        slides = deck_spec.get("slides") if isinstance(deck_spec.get("slides"), list) else []
        for idx, spec in enumerate(slides[:9], start=2):
            if not isinstance(spec, dict):
                continue
            layout = str(spec.get("layout_type") or "two_column").strip().lower()
            chart_path = resolve_chart_path(spec, chart_paths, used_chart_refs)
            chart_spec = resolve_chart_spec(spec, chart_specs)
            if layout in {"big_stat", "big-stat"}:
                render_big_stat_slide(spec, idx)
            elif layout in {"three_column_cards", "three-column-cards", "three_column"}:
                render_three_column_slide(spec, idx)
            elif layout in {"text_and_image", "text-and-image"}:
                render_text_and_image_slide(spec, idx, chart_path, chart_spec)
            elif layout in {"chart", "chart_focus"}:
                render_chart_focus_slide(spec, idx, chart_path, chart_spec)
            elif layout == "hero":
                render_hero_slide(spec, idx)
            elif layout == "thesis":
                render_thesis_slide(spec, idx, chart_path)
            elif layout == "metrics_dashboard":
                render_metrics_slide(spec, idx, chart_path)
            elif layout == "risk_matrix":
                render_risk_matrix_slide(spec, idx, chart_path)
            elif layout == "closing_recommendation":
                render_closing_slide(spec, idx)
            else:
                render_two_column_slide(spec, idx, chart_path)

    if len(prs.slides) < 3:
        render_closing_slide(
            {
                "section_label": "Recommendation",
                "headline": "Committee actions and monitoring plan",
                "takeaway": deck_spec.get("recommendation", ""),
                "bullets": [
                    "Validate assumptions against internal model and consensus.",
                    "Pressure-test downside scenarios before position sizing.",
                    "Define buy, hold, trim, or watchlist decision criteria.",
                ],
            },
            len(prs.slides) + 1,
        )

    # Guardrail: ensure max 10 slides
    while len(prs.slides) > 10:
        r_id = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[-1]

    prs.save(output_path)
    return output_path
